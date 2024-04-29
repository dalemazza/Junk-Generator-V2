import os
import random
import string
import win32com.client as win32
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

def add_junk_to_word_with_com(filename):
    try:
        word = win32.gencache.EnsureDispatch('Word.Application')
        word.Visible = False  # Opens word and inserts junk, set to true for debugging

        # Ensure an absolute path for use of COM
        abs_filename = os.path.abspath(filename)
        doc = word.Documents.Open(abs_filename)

        for _ in range(10):  # Add paragraphs of junk text
            junk_text = ''.join(random.choices(string.ascii_letters + string.digits, k=5000)) # Amount of junk characters
            doc.Content.InsertAfter(junk_text + '\n')

        doc.Save()
        doc.Close()
        word.Quit()
    except Exception as e:
        print(f"An error occurred: {e}")
        if word:
            word.Quit()

def add_junk_to_word(filename):
    if filename.endswith('.docm'):
        add_junk_to_word_with_com(filename)
    else:
        doc = Document(filename)
        for _ in range(10):  # Add paragraphs of junk text
            junk_text = ''.join(random.choices(string.ascii_letters + string.digits, k=5000)) # Amount of junk characters
            doc.add_paragraph(junk_text)
        doc.save(filename)

def add_junk_to_powerpoint(filename):
    ppt = Presentation(filename) 
    
    if len(ppt.slides) == 0: # If presentation has no slides create them and give them content
        title_slide_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Khaos"
        subtitle.text = "Junk Generator V2"

    for slide in ppt.slides:
        if not slide.has_notes_slide:
            slide.notes_slide

        notes_slide = slide.notes_slide # Write junk to Notes portion of slides
        for _ in range(5):
            junk_text = ''.join(random.choices(string.ascii_letters + string.digits, k=5000)) # amount of junk characters
            paragraph = notes_slide.notes_text_frame.add_paragraph()
            paragraph.text = junk_text

    ppt.save(filename)

def add_junk_to_excel(filename):
    workbook = load_workbook(filename, keep_vba=True)
    worksheet = workbook.active
    for i in range(1, 51):
        junk_text = ''.join(random.choices(string.ascii_letters + string.digits, k=5000)) # amount of junk characters
        worksheet[f'A{i}'] = junk_text
    workbook.save(filename)

def main():
    ascii_art = """
     _             _       ____                           _             
    | |_   _ _ __ | | __  / ___| ___ _ __   ___ _ __ __ _| |_ ___  _ __ 
 _  | | | | | '_ \| |/ / | |  _ / _ \ '_ \ / _ \ '__/ _` | __/ _ \| '__|
| |_| | |_| | | | |   <  | |_| |  __/ | | |  __/ | | (_| | || (_) | |   
 \___/ \__,_|_| |_|_|\_\  \____|\___|_| |_|\___|_|  \__,_|\__\___/|_|  V2
 Created by @KhaosShield   // .docm .docx .pptx .pptm .xlsm
    """
    print(ascii_art)

    file_path = input("Enter file path: ")
    if not os.path.exists(file_path):
        print("File not found. Exiting...")
        return

    _, file_extension = os.path.splitext(file_path)
    
    if file_extension.lower() in ['.docx', '.docm']:
        add_junk_to_word(file_path)
    elif file_extension.lower() in ['.pptx', '.pptm']:
        add_junk_to_powerpoint(file_path)
    elif file_extension.lower() == '.xlsm':
        add_junk_to_excel(file_path)
    else:
        print("Incorrect filetype. Only .docx .docm .pptx .pptm .xlsm are supported. Exiting...")
        return

    print(f"SUCCESS! {file_path} is now full of junk.")

if __name__ == "__main__":
    main()
